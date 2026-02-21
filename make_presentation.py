from pptx import Presentation
from pptx.util import Inches, Pt

def create_presentation():
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]

    # Slide 1: Title
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Stegano-Share Secure"
    subtitle.text = "An AES-Encrypted Image Steganography Platform\nB.Tech 3rd Year Project"

    # Slide 2: Introduction
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Introduction & Problem Statement"
    tf = body_shape.text_frame
    tf.text = "The Problem: Digital surveillance makes standard messaging vulnerable."
    p = tf.add_paragraph()
    p.text = "The Solution: Combine Cryptography and Steganography."
    p = tf.add_paragraph()
    p.text = "Goal: Hide AES-encrypted data inside normal-looking images."

    # Slide 3: LSB Algorithm
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "How it Works: LSB Steganography"
    tf = body_shape.text_frame
    tf.text = "Images are grids of pixels (Red, Green, Blue)."
    p = tf.add_paragraph()
    p.text = "Values range from 0-255 (8-bit binary)."
    p = tf.add_paragraph()
    p.text = "Altering the Least Significant Bit (LSB) changes the color invisibly."
    p = tf.add_paragraph()
    p.text = "We replace this bit with binary data from our secret text."

    # Slide 4: Security Layer
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "The Security Layer: AES Encryption"
    tf = body_shape.text_frame
    tf.text = "LSB alone is vulnerable to extraction scripts."
    p = tf.add_paragraph()
    p.text = "Implementation: AES Encryption via Python Cryptography."
    p = tf.add_paragraph()
    p.text = "Process: Text + Password -> AES Ciphertext -> Hidden via LSB."

    # Slide 5: Architecture
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "System Architecture"
    tf = body_shape.text_frame
    tf.text = "Frontend (React.js): Captures image, password, and text payload."
    p = tf.add_paragraph()
    p.text = "Backend (Flask): REST API receives payload and handles image processing."
    p = tf.add_paragraph()
    p.text = "Engine (Python): Executes PBKDF2 key derivation and bitwise operations."

    # Slide 6: Tech Stack
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Technology Stack"
    tf = body_shape.text_frame
    tf.text = "Frontend: React.js, HTML5, CSS3, Axios"
    p = tf.add_paragraph()
    p.text = "Backend: Python 3, Flask"
    p = tf.add_paragraph()
    p.text = "Libraries: Pillow (PIL), Cryptography, Werkzeug"

    # Slide 7: Challenge
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Key Challenge: JPG vs PNG Compression"
    tf = body_shape.text_frame
    tf.text = "Challenge: Encoded messages were destroyed upon saving."
    p = tf.add_paragraph()
    p.text = "Diagnosis: JPEG uses Lossy Compression, overwriting LSB modifications."
    p = tf.add_paragraph()
    p.text = "Solution: Enforced PNG output, which utilizes Lossless Compression."

    # Slide 8: Demo
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Live Demonstration"
    subtitle.text = "Switching to Localhost Environment..."

    # Slide 9: Future Scope
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Future Scope"
    tf = body_shape.text_frame
    tf.text = "Audio/Video Steganography (.mp4, .wav)."
    p = tf.add_paragraph()
    p.text = "Frontend payload capacity checking algorithms."
    p = tf.add_paragraph()
    p.text = "Database integration for secure, one-time shareable links."

    # Slide 10: Q&A
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Thank You"
    subtitle.text = "Questions?"

    # Save the file
    prs.save('Stegano_Share_Presentation.pptx')
    print("Success! 'Stegano_Share_Presentation.pptx' has been created in your folder.")

if __name__ == '__main__':
    create_presentation()